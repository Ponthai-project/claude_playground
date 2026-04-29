from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# Colors
AWS_ORANGE = RGBColor(0xFF, 0x99, 0x00)
AWS_NAVY = RGBColor(0x23, 0x2F, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=AWS_NAVY):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_title_bar(slide, title_text, subtitle_text=None):
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AWS_NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12), Inches(0.8))
    tf = tb.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle_text:
        acc = slide.shapes.add_shape(1, 0, Inches(1.3), Inches(0.08), Inches(6.2))
        acc.fill.solid()
        acc.fill.fore_color.rgb = AWS_ORANGE
        acc.line.fill.background()

def add_textbox(slide, text, left, top, width, height, size=16, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return tb

def add_bullet_box(slide, items, left, top, width, height, title=None, title_color=AWS_ORANGE, text_color=DARK_GRAY, base_size=15, bg_color=None):
    if bg_color:
        box = slide.shapes.add_shape(1, left - Inches(0.1), top - Inches(0.1), width + Inches(0.2), height + Inches(0.2))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.fill.background()

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    if title:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(base_size + 1)
        p.font.bold = True
        p.font.color.rgb = title_color

    for item in items:
        p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, is_orange = item
        else:
            text, is_orange = item, False
        p.text = text
        p.font.size = Pt(base_size)
        p.font.color.rgb = AWS_ORANGE if is_orange else text_color
        p.space_before = Pt(2)


# ===== スライド1: 表紙 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, AWS_NAVY)

accent = slide.shapes.add_shape(1, 0, Inches(2.8), prs.slide_width, Inches(0.06))
accent.fill.solid()
accent.fill.fore_color.rgb = AWS_ORANGE
accent.line.fill.background()
accent2 = slide.shapes.add_shape(1, 0, Inches(5.3), prs.slide_width, Inches(0.06))
accent2.fill.solid()
accent2.fill.fore_color.rgb = AWS_ORANGE
accent2.line.fill.background()

add_textbox(slide, "AWS Certified", Inches(1), Inches(1.2), Inches(11), Inches(0.8), size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "AI Practitioner", Inches(1), Inches(1.9), Inches(11), Inches(1.0), size=48, bold=True, color=AWS_ORANGE, align=PP_ALIGN.CENTER)
add_textbox(slide, "AIF-C01", Inches(1), Inches(2.9), Inches(11), Inches(0.6), size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, "試験対策ガイド ─ 全体像を掴む", Inches(1), Inches(3.6), Inches(11), Inches(0.7), size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "2026年版", Inches(1), Inches(5.5), Inches(11), Inches(0.5), size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ===== スライド2: 試験概要 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "試験概要")

rows = [
    ("受験料", "$100（約16,500円）"),
    ("試験時間", "90分"),
    ("問題数", "65問（採点対象50問＋ノースコア15問）"),
    ("合格ライン", "700点 / 1,000点満点"),
    ("レベル", "Foundational（最も基礎的なレベル）"),
    ("前提条件", "なし（経験6ヶ月未満でも可）"),
    ("対象者", "エンジニア・非エンジニア（PM・ビジネス職）問わず"),
    ("出題形式", "単一選択／複数選択／順序並べ替え"),
]

for i, (key, val) in enumerate(rows):
    top = Inches(1.55) + i * Inches(0.68)
    bg = slide.shapes.add_shape(1, Inches(0.4), top - Inches(0.05), Inches(12.5), Inches(0.6))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(0xEE, 0xEE, 0xEE)
    bg.line.fill.background()

    ktb = slide.shapes.add_textbox(Inches(0.5), top, Inches(3.2), Inches(0.55))
    ktf = ktb.text_frame
    kp = ktf.add_paragraph()
    kp.text = key
    kp.font.size = Pt(16)
    kp.font.bold = True
    kp.font.color.rgb = AWS_NAVY

    vtb = slide.shapes.add_textbox(Inches(3.9), top, Inches(9), Inches(0.55))
    vtf = vtb.text_frame
    vtf.word_wrap = True
    vp = vtf.add_paragraph()
    vp.text = val
    vp.font.size = Pt(16)
    vp.font.color.rgb = DARK_GRAY


# ===== スライド3: 出題ドメインと配点 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "出題ドメインと配点（5ドメイン）")

domains = [
    ("Domain 1", "AI/MLの基礎", 20, AWS_NAVY),
    ("Domain 2", "生成AIの基礎", 24, RGBColor(0x1A, 0x5E, 0x9A)),
    ("Domain 3", "基盤モデルの活用", 28, RGBColor(0x0D, 0x7E, 0x6E)),
    ("Domain 4", "責任あるAI", 14, RGBColor(0x8B, 0x2F, 0x00)),
    ("Domain 5", "セキュリティ・コンプライアンス・ガバナンス", 14, RGBColor(0x6B, 0x21, 0x6B)),
]

bar_left = Inches(5.2)
bar_max_width = Inches(7.5)
bar_height = Inches(0.65)
total_pct = 28  # max for scaling

for i, (code, name, pct, color) in enumerate(domains):
    top = Inches(1.6) + i * Inches(1.0)

    label_tb = slide.shapes.add_textbox(Inches(0.4), top + Inches(0.1), Inches(4.6), Inches(0.5))
    tf = label_tb.text_frame
    p = tf.add_paragraph()
    p.text = f"{code}: {name}"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    bar_bg = slide.shapes.add_shape(1, bar_left, top, bar_max_width, bar_height)
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    bar_bg.line.fill.background()

    bar_w = bar_max_width * pct / total_pct
    bar = slide.shapes.add_shape(1, bar_left, top, int(bar_w), bar_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    pct_tb = slide.shapes.add_textbox(bar_left + int(bar_w) + Inches(0.1), top + Inches(0.1), Inches(0.8), Inches(0.5))
    tf2 = pct_tb.text_frame
    p2 = tf2.add_paragraph()
    p2.text = f"{pct}%"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = color


# ===== スライド4: ドメイン1 AI/MLの基礎 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "Domain 1：AI / ML の基礎  ─ 20%")

add_bullet_box(slide, [
    "● 機械学習の種類：教師あり／なし／強化学習",
    "● 主要アルゴリズム：回帰・分類・クラスタリング・NN",
    "● モデル評価指標：精度・適合率・再現率・F1・RMSE",
    "● 過学習・過少学習・バイアス・バリアンス",
    "● データ前処理：正規化・欠損値処理・特徴量エンジニアリング",
], Inches(0.4), Inches(1.6), Inches(6.0), Inches(3.5), title="MLの基本概念", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● Amazon SageMaker（モデル開発・訓練・デプロイ）",
    "● Amazon Rekognition（画像・動画認識）",
    "● Amazon Comprehend（テキスト分析・NLP）",
    "● Amazon Forecast（時系列予測）",
    "● Amazon Personalize（レコメンデーション）",
    "● Amazon Textract（OCR・ドキュメント解析）",
], Inches(6.6), Inches(1.6), Inches(6.3), Inches(3.5), title="主要AWSサービス", text_color=DARK_GRAY, bg_color=WHITE)

tip = slide.shapes.add_shape(1, Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.7))
tip.fill.solid()
tip.fill.fore_color.rgb = AWS_NAVY
tip.line.fill.background()
add_textbox(slide, "★ ポイント：各アルゴリズムの「どんな問題に使うか」を整理すると得点しやすい", Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.6), size=14, color=WHITE)


# ===== スライド5: ドメイン2 生成AIの基礎 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "Domain 2：生成AI の基礎  ─ 24%")

add_bullet_box(slide, [
    "● 大規模言語モデル（LLM）の仕組み",
    "● Transformer・アテンション機構",
    "● トークン・コンテキストウィンドウ",
    "● プロンプトエンジニアリング（Zero-shot / Few-shot）",
    "● ファインチューニング vs RAG（検索拡張生成）",
    "● ハルシネーション（幻覚）とその対策",
    "● 拡散モデル・GANによる画像生成",
], Inches(0.4), Inches(1.6), Inches(6.0), Inches(4.2), title="生成AI基本概念", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● Amazon Bedrock（基盤モデルAPIサービス）",
    "● Amazon Titan（AWS独自モデル群）",
    "● Anthropic Claude / Meta Llama（Bedrock経由）",
    "● Amazon Q（ビジネス向けAIアシスタント）",
    "● Amazon CodeWhisperer（コード補完）",
    "● Stable Diffusion（画像生成）",
], Inches(6.6), Inches(1.6), Inches(6.3), Inches(3.5), title="主要AWSサービス", text_color=DARK_GRAY, bg_color=WHITE)

tip = slide.shapes.add_shape(1, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.7))
tip.fill.solid()
tip.fill.fore_color.rgb = AWS_NAVY
tip.line.fill.background()
add_textbox(slide, "★ ポイント：RAGとファインチューニングの使い分け・Bedrockの役割を押さえる", Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.6), size=14, color=WHITE)


# ===== スライド6: ドメイン3 基盤モデルの活用 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "Domain 3：基盤モデルの活用  ─ 28%（最重要）")

add_bullet_box(slide, [
    "● Foundation Model（FM）の選び方・評価",
    "● モデルカスタマイズ：Fine-tuning / Continued Pre-training",
    "● RAG（Retrieval-Augmented Generation）構築",
    "● ベクトルデータベースの活用（埋め込みベクトル）",
    "● エージェント型AI・ツール使用",
    "● モデル評価：BLEU・ROUGE・人間評価",
    "● コスト最適化（モデル選択・推論コスト）",
], Inches(0.4), Inches(1.6), Inches(6.0), Inches(4.2), title="活用パターンと概念", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● Amazon Bedrock（FMのAPI呼び出し・管理）",
    "● Bedrock Knowledge Bases（RAG構築）",
    "● Bedrock Agents（エージェント構築）",
    "● Bedrock Guardrails（安全フィルタ）",
    "● Amazon OpenSearch（ベクトル検索）",
    "● Amazon Kendra（エンタープライズ検索）",
    "● AWS Lambda（推論エンドポイント連携）",
], Inches(6.6), Inches(1.6), Inches(6.3), Inches(4.2), title="主要AWSサービス", text_color=DARK_GRAY, bg_color=WHITE)

tip = slide.shapes.add_shape(1, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.7))
tip.fill.solid()
tip.fill.fore_color.rgb = AWS_ORANGE
tip.line.fill.background()
add_textbox(slide, "★★ 最重要ドメイン（28%）！ Bedrock関連サービスを完全に理解すること", Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.6), size=14, color=AWS_NAVY, bold=True)


# ===== スライド7: ドメイン4 責任あるAI =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "Domain 4：責任ある AI  ─ 14%")

add_bullet_box(slide, [
    "● 公平性（Fairness）：バイアスの特定と軽減",
    "● 透明性（Transparency）：意思決定の説明可能性",
    "● 説明可能性（Explainability）：XAI手法",
    "● プライバシー保護：データ匿名化・差分プライバシー",
    "● 堅牢性（Robustness）：敵対的攻撃への耐性",
    "● ヒューマン・イン・ザ・ループ",
    "● AIガバナンスフレームワーク",
], Inches(0.4), Inches(1.6), Inches(6.0), Inches(4.2), title="責任あるAIの6原則", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● Amazon SageMaker Clarify（バイアス検出）",
    "● Amazon SageMaker Model Monitor（モデル監視）",
    "● Bedrock Guardrails（有害コンテンツフィルタ）",
    "● AWS AI Service Cards（透明性ドキュメント）",
    "● AWS Responsible AI Policy",
], Inches(6.6), Inches(1.6), Inches(6.3), Inches(3.0), title="主要AWSサービス・ツール", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● モデルカードの作成・公開",
    "● バイアス監査の定期実施",
    "● ステークホルダーへの説明責任",
], Inches(6.6), Inches(4.8), Inches(6.3), Inches(1.8), title="ベストプラクティス", text_color=DARK_GRAY, bg_color=WHITE)

tip = slide.shapes.add_shape(1, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.7))
tip.fill.solid()
tip.fill.fore_color.rgb = AWS_NAVY
tip.line.fill.background()
add_textbox(slide, "★ ポイント：各原則の定義と、それに対応するAWSサービスをセットで覚える", Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.6), size=14, color=WHITE)


# ===== スライド8: ドメイン5 セキュリティ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "Domain 5：セキュリティ・コンプライアンス・ガバナンス  ─ 14%")

add_bullet_box(slide, [
    "● IAM（最小権限の原則・ロールベースアクセス）",
    "● データ暗号化（保存時・転送時）",
    "● VPC・プライベートエンドポイント",
    "● 監査ログ・CloudTrail",
    "● インシデント対応・脅威検出",
], Inches(0.4), Inches(1.6), Inches(4.0), Inches(3.5), title="セキュリティ基礎", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● GDPR・HIPAA・SOC 2・ISO 27001",
    "● データ主権・残留規制",
    "● AWS Artifact（コンプライアンスレポート）",
    "● AWS Config（設定コンプライアンス）",
    "● AWS Audit Manager",
], Inches(4.6), Inches(1.6), Inches(4.0), Inches(3.5), title="コンプライアンス", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● AIガバナンスポリシーの策定",
    "● モデルライフサイクル管理",
    "● リスクアセスメント",
    "● AWS Organizations・SCP",
    "● AWS Well-Architected Framework（ML Lens）",
], Inches(9.0), Inches(1.6), Inches(4.0), Inches(3.5), title="ガバナンス", text_color=DARK_GRAY, bg_color=WHITE)

tip = slide.shapes.add_shape(1, Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.7))
tip.fill.solid()
tip.fill.fore_color.rgb = AWS_NAVY
tip.line.fill.background()
add_textbox(slide, "★ ポイント：AWSの共有責任モデル（お客様 vs AWS）のAI文脈での理解が重要", Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.6), size=14, color=WHITE)


# ===== スライド9: 受験方法 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "受験方法")

add_bullet_box(slide, [
    "① aws.amazon.com/certification にアクセス",
    "② AWSアカウント（またはAWS Certificationアカウント）を作成",
    "③「試験を予約する」→ AIF-C01 を選択",
    "④ 受験方式を選択（オンラインまたはテストセンター）",
    "⑤ 日時・会場を指定して支払い完了",
], Inches(0.4), Inches(1.55), Inches(12.5), Inches(2.5), title="申込手順", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● 自宅のPC・Webカメラ・マイクが必要",
    "● Pearson VUE OnVUE ソフトをインストール",
    "● 受験中は画面・周囲を監視員が監視",
    "● 静かな個室環境が必須",
], Inches(0.4), Inches(4.3), Inches(5.9), Inches(2.5), title="オンライン受験（自宅）", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● 全国各地のピアソンVUEテストセンター",
    "● 身分証明書（写真付き）2点が必要",
    "● 持ち込み不可（ノートPC・メモ等）",
    "● 安心して受けたい方におすすめ",
], Inches(6.7), Inches(4.3), Inches(6.2), Inches(2.5), title="テストセンター受験", text_color=DARK_GRAY, bg_color=WHITE)


# ===== スライド10: 学習リソース =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "推奨学習リソース")

add_bullet_box(slide, [
    "● AWS公式試験ガイド（無料PDF）",
    "● AWS Skill Builder（公式学習プラットフォーム）",
    "  - 「AWS Certified AI Practitioner」公式コース",
    "  - 公式模擬試験（20問・無料）",
    "● AWS Black Belt Online Seminar（YouTube）",
    "● AWS公式ドキュメント（Bedrock等）",
], Inches(0.4), Inches(1.55), Inches(5.9), Inches(3.8), title="AWS公式（無料）", text_color=DARK_GRAY, bg_color=WHITE)

add_bullet_box(slide, [
    "● Tutorials Dojo：高品質模擬問題集",
    "● Whizlabs：解説が丁寧・日本語対応あり",
    "● Udemy：動画コース（セール時＄10〜）",
    "● DX/AI研究所ブログ（日本語解説）",
    "● Zenn・Qiita（合格体験記多数）",
], Inches(6.7), Inches(1.55), Inches(6.2), Inches(3.8), title="サードパーティ（有料・無料）", text_color=DARK_GRAY, bg_color=WHITE)

tip = slide.shapes.add_shape(1, Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.7))
tip.fill.solid()
tip.fill.fore_color.rgb = AWS_ORANGE
tip.line.fill.background()
add_textbox(slide, "★ まずはAWS Skill Builderの公式コース＋公式模擬試験（無料）から始めるのが最短ルート", Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.6), size=14, color=AWS_NAVY, bold=True)


# ===== スライド11: 学習ロードマップ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_title_bar(slide, "学習ロードマップ（目安：1〜2ヶ月）")

steps = [
    ("Week 1-2", "全体把握", "試験ガイド読了 → Skill Builderで概要動画 → 各ドメインの用語リスト作成", RGBColor(0x1A, 0x5E, 0x9A)),
    ("Week 3-4", "インプット強化", "ドメイン別に深掘り（特にDomain 2・3） → AWS公式ドキュメントでBedrockを重点学習", RGBColor(0x0D, 0x7E, 0x6E)),
    ("Week 5-6", "問題演習", "Tutorials Dojoで模擬試験 → 間違えた問題の解説を精読 → 弱点ドメインを再インプット", RGBColor(0x8B, 0x2F, 0x00)),
    ("Week 7-8", "仕上げ", "全ドメイン模擬試験を繰り返し → 700点安定 → 本番予約・受験", AWS_NAVY),
]

for i, (period, phase, desc, color) in enumerate(steps):
    top = Inches(1.55) + i * Inches(1.35)

    badge = slide.shapes.add_shape(1, Inches(0.4), top, Inches(1.5), Inches(1.1))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    btb = slide.shapes.add_textbox(Inches(0.4), top + Inches(0.05), Inches(1.5), Inches(1.0))
    btf = btb.text_frame
    btf.word_wrap = True
    bp1 = btf.add_paragraph()
    bp1.text = period
    bp1.font.size = Pt(13)
    bp1.font.bold = True
    bp1.font.color.rgb = WHITE
    bp1.alignment = PP_ALIGN.CENTER
    bp2 = btf.add_paragraph()
    bp2.text = phase
    bp2.font.size = Pt(12)
    bp2.font.color.rgb = AWS_ORANGE
    bp2.alignment = PP_ALIGN.CENTER

    content_bg = slide.shapes.add_shape(1, Inches(2.1), top, Inches(11.0), Inches(1.1))
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = WHITE
    content_bg.line.fill.background()

    dtb = slide.shapes.add_textbox(Inches(2.2), top + Inches(0.2), Inches(10.8), Inches(0.8))
    dtf = dtb.text_frame
    dtf.word_wrap = True
    dp = dtf.add_paragraph()
    dp.text = desc
    dp.font.size = Pt(14)
    dp.font.color.rgb = DARK_GRAY


# ===== スライド12: まとめ・合格のポイント =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, AWS_NAVY)

accent_bar = slide.shapes.add_shape(1, 0, Inches(1.1), prs.slide_width, Inches(0.06))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = AWS_ORANGE
accent_bar.line.fill.background()

add_textbox(slide, "まとめ・合格のポイント", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), size=28, bold=True, color=WHITE)

add_bullet_box(slide, [
    "✅  Domain 3（基盤モデルの活用・28%）を最優先で攻略する",
    "✅  Amazon Bedrockとその関連サービスを完全に理解する",
    "✅  RAG vs ファインチューニングの使い分けを説明できるようにする",
    "✅  責任あるAIの6原則と対応AWSサービスをセットで暗記する",
    "✅  Tutorials Dojoなどで模擬試験を繰り返し、700点を安定させる",
    "✅  AWS Skill Builderの公式模擬試験（無料）は必ず受ける",
], Inches(0.5), Inches(1.3), Inches(12.4), Inches(3.8), text_color=WHITE, bg_color=None, base_size=16)

bottom_bar = slide.shapes.add_shape(1, 0, Inches(5.8), prs.slide_width, Inches(1.7))
bottom_bar.fill.solid()
bottom_bar.fill.fore_color.rgb = AWS_ORANGE
bottom_bar.line.fill.background()

add_textbox(slide, "Foundationalレベルは基礎・概念理解が中心。コーディング不要。", Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.6), size=18, bold=True, color=AWS_NAVY, align=PP_ALIGN.CENTER)
add_textbox(slide, "1〜2ヶ月の学習で十分合格可能。まずAWS Skill Builderから始めよう！", Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6), size=16, color=AWS_NAVY, align=PP_ALIGN.CENTER)


# ===== スライド13: 参考文献・公式リソース =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

bar13 = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
bar13.fill.solid()
bar13.fill.fore_color.rgb = AWS_NAVY
bar13.line.fill.background()

acc13 = slide.shapes.add_shape(1, 0, Inches(1.0), prs.slide_width, Inches(0.05))
acc13.fill.solid()
acc13.fill.fore_color.rgb = AWS_ORANGE
acc13.line.fill.background()

add_textbox(slide, "参考文献・公式リソース", Inches(0.4), Inches(0.13), Inches(12), Inches(0.75), size=26, bold=True, color=WHITE)

LINK_BLUE = RGBColor(0x00, 0x56, 0xB3)

def ref_section(slide, title, x, y, w):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.35))
    tf = tb.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE
    return y + Inches(0.37)

def ref_item(slide, label, url, x, y, w):
    item_h = Inches(0.72)
    box = slide.shapes.add_shape(1, x, y, w, item_h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    ltb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.05), w - Inches(0.24), Inches(0.3))
    ltf = ltb.text_frame
    ltf.word_wrap = True
    lp = ltf.add_paragraph()
    lp.text = label
    lp.font.size = Pt(12)
    lp.font.bold = True
    lp.font.color.rgb = AWS_NAVY

    utb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.38), w - Inches(0.24), Inches(0.28))
    utf = utb.text_frame
    utf.word_wrap = True
    up = utf.add_paragraph()
    up.text = url
    up.font.size = Pt(8.5)
    up.font.color.rgb = LINK_BLUE
    up.font.underline = True

    return y + item_h + Inches(0.05)

col_w = Inches(6.1)
x_l = Inches(0.35)
x_r = Inches(6.88)
y_l = Inches(1.15)
y_r = Inches(1.15)

y_l = ref_section(slide, "【試験公式】", x_l, y_l, col_w)
y_l = ref_item(slide, "試験公式ページ", "https://aws.amazon.com/certification/certified-ai-practitioner/", x_l, y_l, col_w)
y_l = ref_item(slide, "試験ガイド（日本語PDF）", "https://d1.awsstatic.com/ja_JP/training-and-certification/docs-ai-practitioner/AWS-Certified-AI-Practitioner_Exam-Guide.pdf", x_l, y_l, col_w)
y_l = ref_item(slide, "試験ガイド（AWS Docs）", "https://docs.aws.amazon.com/aws-certification/latest/examguides/ai-practitioner-01.html", x_l, y_l, col_w)
y_l = ref_item(slide, "試験予約（Pearson VUE）", "https://aws.amazon.com/certification/certification-prep/testing/", x_l, y_l, col_w)
y_l = ref_item(slide, "AWS認定 FAQ", "https://aws.amazon.com/certification/faqs/", x_l, y_l, col_w)

y_r = ref_section(slide, "【学習リソース（公式）】", x_r, y_r, col_w)
y_r = ref_item(slide, "AWS Skill Builder - 試験準備コース（無料）", "https://skillbuilder.aws/learn/X83W99WJXA/exam-prep-standard-course-aws-certified-ai-practitioner-aif-c01/KUW4WB2K4B", x_r, y_r, col_w)
y_r = ref_item(slide, "AWS Skill Builder - 公式模擬問題集（無料）", "https://explore.skillbuilder.aws/learn/course/19790/exam-prep-official-practice-question-set-aws-certified-ai-practitioner-aif-c01-english", x_r, y_r, col_w)
y_r = ref_item(slide, "AWS Skill Builder - 試験準備プラン", "https://skillbuilder.aws/category/exam-prep/ai-practitioner", x_r, y_r, col_w)

y_r += Inches(0.12)
y_r = ref_section(slide, "【AWSサービス公式ドキュメント】", x_r, y_r, col_w)
y_r = ref_item(slide, "Amazon Bedrock 公式ドキュメント", "https://docs.aws.amazon.com/bedrock/latest/userguide/what-is-bedrock.html", x_r, y_r, col_w)
y_r = ref_item(slide, "Amazon SageMaker 公式ドキュメント", "https://docs.aws.amazon.com/sagemaker/", x_r, y_r, col_w)


# 保存
import os
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AWS_AI_Practitioner_Overview.pptx")
prs.save(output_path)
print(f"保存完了: {output_path}")
