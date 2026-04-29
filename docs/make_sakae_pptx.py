from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- 色定数（白ベース） ----
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GOLD       = RGBColor(0xA0, 0x78, 0x20)
GOLD_DARK  = RGBColor(0x78, 0x58, 0x10)
GOLD_LIGHT = RGBColor(0xE8, 0xD4, 0x8A)
GOLD_BG    = RGBColor(0xFD, 0xF8, 0xED)
DARK       = RGBColor(0x1C, 0x1C, 0x1C)
CHARCOAL   = RGBColor(0x3C, 0x3C, 0x3C)
MID_GRAY   = RGBColor(0x70, 0x70, 0x70)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
RED_ACCENT = RGBColor(0x8B, 0x1A, 0x1A)
CREAM_BG   = RGBColor(0xF5, 0xF0, 0xE8)
PANEL_BG   = RGBColor(0xF0, 0xEB, 0xE0)

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, x, y, w, h, fill, line_color=None, line_w=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        if line_w: s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h, size=16, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.italic = italic; p.font.color.rgb = color; p.alignment = align
    return tb

def header(slide, title, subtitle):
    add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
    add_rect(slide, 0, 0, 13.33, 1.35, GOLD_BG)
    add_rect(slide, 0, 1.35, 13.33, 0.05, GOLD)
    add_rect(slide, 0.4, 0.15, 0.07, 1.05, GOLD)
    add_text(slide, title,    0.65, 0.12, 11, 0.85, size=36, bold=True, color=GOLD_DARK)
    add_text(slide, subtitle, 0.65, 0.95, 10, 0.38, size=12, color=GOLD, italic=True)

# ============================================================
# 1. 表紙
# ============================================================
def slide_cover(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
    add_rect(slide, 0, 0, 13.33, 0.28, GOLD)
    add_rect(slide, 0, 7.22, 13.33, 0.28, GOLD)
    add_rect(slide, 0.35, 0.45, 0.07, 6.6, GOLD)
    add_rect(slide, 0.56, 0.45, 0.02, 6.6, GOLD_LIGHT)
    add_rect(slide, 12.72, 0.45, 0.07, 6.6, GOLD)
    add_rect(slide, 1.2, 1.5, 10.9, 4.3, GOLD_BG)
    add_rect(slide, 1.2, 1.5, 10.9, 0.07, GOLD)
    add_rect(slide, 1.2, 5.73, 10.9, 0.07, GOLD)
    add_text(slide, "上野 榮（さかえ）", 1.3, 1.72, 10.7, 1.4,
             size=56, bold=True, color=GOLD_DARK, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.5, 3.1, 6.3, 0.05, GOLD_LIGHT)
    add_text(slide, "完全ガイド", 1.3, 3.22, 10.7, 1.0,
             size=38, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, "— 江戸前鮨の本格派 —", 1.3, 4.28, 10.7, 0.72,
             size=20, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
    add_text(slide, "東京都台東区上野　c-roadビル12F　｜　銀座久兵衛の系譜",
             1.3, 5.05, 10.7, 0.55, size=14, color=CHARCOAL, align=PP_ALIGN.CENTER)
    add_text(slide, "調査資料", 0.9, 6.68, 3, 0.45, size=13, color=MID_GRAY)

# ============================================================
# 2. お店概要
# ============================================================
def slide_overview(prs):
    slide = blank_slide(prs)
    header(slide, "お店概要", "STORE INFORMATION")
    # 左カード
    add_rect(slide, 0.4, 1.5, 5.9, 5.65, CREAM_BG)
    add_rect(slide, 0.4, 1.5, 0.09, 5.65, GOLD)
    add_text(slide, "基本情報", 0.67, 1.58, 5.4, 0.6, size=18, bold=True, color=GOLD_DARK)
    add_rect(slide, 0.4, 2.18, 5.9, 0.04, GOLD_LIGHT)
    info = [
        ("住所",    "東京都台東区上野4-4-5\nc-roadビル12F"),
        ("電話",    "03-6284-4731"),
        ("ランチ",  "12:00 〜 14:30"),
        ("ディナー","17:00 〜 23:00（最終入店 21:00）"),
        ("定休日",  "不定休（年末年始等休業あり）"),
        ("個室料",  "個室チャージ・サービス料なし"),
    ]
    y = 2.28
    for label, val in info:
        add_text(slide, label, 0.67, y, 1.5, 0.55, size=13, bold=True, color=GOLD)
        add_text(slide, val,   2.25, y, 3.85, 0.65, size=14, color=DARK, wrap=True)
        add_rect(slide, 0.5, y+0.72, 5.7, 0.02, LIGHT_GRAY)
        y += 0.82
    # 右カード
    add_rect(slide, 6.9, 1.5, 6.05, 5.65, CREAM_BG)
    add_rect(slide, 6.9, 1.5, 0.09, 5.65, GOLD)
    add_text(slide, "アクセス", 7.17, 1.58, 5.6, 0.6, size=18, bold=True, color=GOLD_DARK)
    add_rect(slide, 6.9, 2.18, 6.05, 0.04, GOLD_LIGHT)
    access = [
        ("🚇", "上野御徒町駅（A5出口）", "徒歩 0分"),
        ("🚇", "上野広小路駅",           "徒歩 2分"),
        ("🚇", "御徒町駅",               "徒歩 3分"),
        ("🚇", "上野駅（公園口）",       "徒歩 7分"),
    ]
    y = 2.28
    for icon, sta, time in access:
        add_text(slide, f"{icon}  {sta}", 7.17, y, 4.0, 0.55, size=15, color=DARK)
        add_text(slide, time, 10.8, y, 1.9, 0.55, size=15, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
        add_rect(slide, 7.0, y+0.6, 5.8, 0.02, LIGHT_GRAY)
        y += 0.74
    # 評価バッジ
    add_rect(slide, 6.9, 5.25, 6.05, 1.9, GOLD_BG)
    add_rect(slide, 6.9, 5.25, 6.05, 0.06, GOLD)
    add_text(slide, "食べログ評価", 7.1, 5.35, 5.7, 0.5, size=15, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "3.56", 7.1, 5.72, 5.7, 1.0, size=54, bold=True, color=GOLD_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, "534件のレビュー", 7.1, 6.78, 5.7, 0.35, size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# 3. コンセプト
# ============================================================
def slide_concept(prs):
    slide = blank_slide(prs)
    header(slide, "コンセプト", "CONCEPT & PHILOSOPHY")
    add_rect(slide, 0.4, 1.5, 12.53, 1.0, GOLD_BG)
    add_rect(slide, 0.4, 1.5, 12.53, 0.06, GOLD)
    add_rect(slide, 0.4, 2.44, 12.53, 0.06, GOLD)
    add_text(slide, "「 銀座久兵衛の系譜を継ぐ、正統派江戸前鮨 」",
             0.5, 1.6, 12.3, 0.78, size=22, bold=True, color=GOLD_DARK,
             align=PP_ALIGN.CENTER, italic=True)
    points = [
        ("江戸前鮨の哲学",
         "旬の素材を主役に、シャリとネタの対話を大切にする伝統の技法。\n"
         "その日の仕入れで最高のネタを選び抜き、丁寧に仕込む。"),
        ("久兵衛の系譜",
         "銀座を代表する名店「久兵衛」の流れを汲む正統派スタイル。\n"
         "格式と品格を守りながら、上野という地で独自の花を開かせた。"),
        ("素材への真摯な姿勢",
         "大将自ら毎日市場へ出向き、食材の個性を見極める。\n"
         "四季折々の旬を、最もおいしい瞬間に提供することにこだわる。"),
    ]
    y = 2.6
    for i, (title, body) in enumerate(points):
        add_rect(slide, 0.4, y, 12.53, 1.58, PANEL_BG if i%2==0 else WHITE)
        add_rect(slide, 0.4, y, 0.09, 1.58, GOLD)
        add_rect(slide, 0.4, y+1.54, 12.53, 0.04, LIGHT_GRAY)
        add_text(slide, title, 0.67, y+0.12, 12.0, 0.58, size=18, bold=True, color=GOLD_DARK)
        add_text(slide, body,  0.67, y+0.68, 12.0, 0.85, size=14, color=DARK, wrap=True)
        y += 1.62

# ============================================================
# 4. 店内の雰囲気
# ============================================================
def slide_atmosphere(prs):
    slide = blank_slide(prs)
    header(slide, "店内の雰囲気", "ATMOSPHERE & INTERIOR")
    features = [
        ("白木のカウンター",
         "職人の手技が間近で見られる特等席。大将と会話しながら、その日のおすすめを聞ける贅沢な空間。",
         GOLD),
        ("格子の半個室テーブル席",
         "格子で仕切られたプライベート感のある席。大切な人との食事や接待にも最適。個室チャージ無し。",
         RGBColor(0x4A, 0x78, 0x4A)),
        ("12階からの眺望",
         "上野の街を一望できる高層階。都会の喧騒を離れた、非日常的な静寂の中で鮨を楽しめる。",
         RGBColor(0x3A, 0x5A, 0x8A)),
        ("サービス料なし",
         "個室チャージ・サービス料は一切不要。本物の味に集中できる良心的な価格設定が好評。",
         RED_ACCENT),
    ]
    positions = [(0.4, 1.5), (6.9, 1.5), (0.4, 4.55), (6.9, 4.55)]
    for (x, y), (title, body, accent) in zip(positions, features):
        add_rect(slide, x, y, 6.05, 2.8, PANEL_BG)
        add_rect(slide, x, y, 0.09, 2.8, accent)
        add_rect(slide, x, y+2.76, 6.05, 0.04, LIGHT_GRAY)
        add_text(slide, title, x+0.28, y+0.2,  5.55, 0.62, size=18, bold=True, color=GOLD_DARK)
        add_rect(slide, x+0.28, y+0.82, 5.6, 0.04, GOLD_LIGHT)
        add_text(slide, body,  x+0.28, y+0.96, 5.6, 1.65, size=13, color=DARK, wrap=True)

# ============================================================
# 5. コース・料金一覧
# ============================================================
def slide_courses(prs):
    slide = blank_slide(prs)
    header(slide, "コース・料金一覧", "COURSE MENU & PRICING")
    courses = [
        ("榮 江戸前寿司握り",    "6,000円",         "握りのみを楽しむシンプルコース",            "L/D"),
        ("ランチ おまかせ上",    "7,000円",         "昼だけの特別おまかせ",                      "L"),
        ("おためしコース",       "9,000円",         "初めての方にもおすすめの入門コース",          "D"),
        ("おまかせコース",       "12,000〜13,200円", "つまみと握りのバランスが絶妙な本格コース",   "D"),
        ("緑コース（縁）",       "16,500円",        "季節の食材をふんだんに使った上質なおまかせ",  "D"),
        ("福20000コース",        "20,000円",        "最高峰の食材で仕立てる、究極のおまかせ",      "D"),
    ]
    col_x = [0.4, 3.55, 6.1, 11.65]
    col_w = [3.05, 2.45, 5.45, 1.25]
    hdrs  = ["コース名", "価格（税込）", "特徴", "時間"]
    add_rect(slide, 0.4, 1.5, 12.53, 0.65, GOLD)
    for h, cx, cw in zip(hdrs, col_x, col_w):
        al = PP_ALIGN.CENTER if h in ("価格（税込）","時間") else PP_ALIGN.LEFT
        add_text(slide, h, cx+0.1, 1.55, cw, 0.52, size=14, bold=True, color=WHITE, align=al)
    y = 2.15
    for j, (name, price, desc, timing) in enumerate(courses):
        bg = PANEL_BG if j%2==0 else WHITE
        if j == len(courses)-1: bg = GOLD_BG
        add_rect(slide, 0.4, y, 12.53, 0.88, bg)
        add_rect(slide, 0.4, y+0.85, 12.53, 0.03, LIGHT_GRAY)
        data   = [name, price, desc, timing]
        aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER]
        colors = [DARK, GOLD_DARK if j==len(courses)-1 else GOLD, CHARCOAL, MID_GRAY]
        bolds  = [True, True, False, False]
        for cx, cw, d, al, col, bo in zip(col_x, col_w, data, aligns, colors, bolds):
            add_text(slide, d, cx+0.1, y+0.1, cw, 0.68, size=13, bold=bo, color=col, align=al, wrap=True)
        y += 0.9
    add_text(slide,
             "L=ランチ　D=ディナー　※料金は税込・サービス料なし。内容は仕入れ状況により変わります。",
             0.4, 7.08, 12.5, 0.38, size=11, color=MID_GRAY)

# ============================================================
# 6. おまかせコースの流れ
# ============================================================
def slide_flow(prs):
    slide = blank_slide(prs)
    header(slide, "おまかせコースの流れ", "OMAKASE COURSE FLOW  ※ある日の一例")
    steps = [
        ("01", "先付け",           "季節の前菜で食欲を誘う"),
        ("02", "握り（前半）",     "中トロ・のどぐろ・シロイカなど4カン"),
        ("03", "焼き物",           "蒸し鮑の肝ソース和えなど"),
        ("04", "子丼",             "いくら雲丹の贅沢なひと椀"),
        ("05", "握り（中盤）",     "馬糞雲丹・赤身漬け・車海老など4カン"),
        ("06", "小鉢 + 茶碗蒸し", "繊細な仕事が光る箸休め"),
        ("07", "握り（後半）",     "金目鯛・炙り大トロ・穴子など4カン"),
        ("08", "巻物・玉子・止椀", "コースの締めくくり"),
        ("09", "デザート",         "季節のひと品で余韻を楽しむ"),
    ]
    cw, ch, pad = 4.1, 1.95, 0.07
    for i, (num, step, desc) in enumerate(steps):
        col = i // 3; row = i % 3
        x = 0.4  + col * (cw + pad)
        y = 1.45 + row * (ch + pad)
        add_rect(slide, x, y, cw, ch, PANEL_BG if col%2==0 else WHITE)
        add_rect(slide, x, y, 0.09, ch, GOLD)
        add_rect(slide, x, y+ch-0.04, cw, 0.04, LIGHT_GRAY)
        add_text(slide, num,  x+0.22, y+0.14, 0.75, 0.6,      size=22, bold=True, color=GOLD)
        add_text(slide, step, x+1.02, y+0.17, cw-1.15, 0.6,   size=17, bold=True, color=GOLD_DARK)
        add_text(slide, desc, x+0.22, y+0.98, cw-0.38, 0.82,  size=13, color=CHARCOAL, wrap=True)

# ============================================================
# 7. ネタ一覧
# ============================================================
def slide_neta(prs):
    slide = blank_slide(prs)
    header(slide, "代表的なネタ", "SIGNATURE INGREDIENTS")
    neta = [
        ("中トロ",      "脂の旨味と赤身のコクの絶妙なバランス",  GOLD),
        ("炙り大トロ",  "炭塩で仕上げる贅沢な一カン",            GOLD),
        ("のどぐろ",    "山陰の高級魚。脂の乗りが抜群",           RGBColor(0xB0,0x60,0x20)),
        ("赤身漬け",    "江戸前の真髄。醤油漬けで旨味を凝縮",    RED_ACCENT),
        ("シロイカ",    "竹炭塩で仕上げる透明感ある甘み",         RGBColor(0x50,0x70,0x90)),
        ("蒸し鮑",      "肝ソース和え。職人技の結晶",             RGBColor(0x40,0x70,0x50)),
        ("馬糞雲丹",    "濃厚な甘みと磯の香りが広がる",           RGBColor(0xA0,0x80,0x10)),
        ("車海老",      "海老味噌添え。重層的な旨味",             RGBColor(0xA0,0x40,0x20)),
        ("かんぱち",    "コリコリ食感と淡白な旨味",               RGBColor(0x30,0x60,0x90)),
        ("金目鯛",      "皮目の香ばしさと上品な脂",               RGBColor(0xC0,0x50,0x20)),
        ("穴子",        "ふっくら柔らか。ツメとの相性が抜群",     RGBColor(0x70,0x50,0x20)),
        ("いくら雲丹丼","子丼として供される贅沢な一品",            RGBColor(0x80,0x20,0x20)),
    ]
    cw, ch, gap = 3.08, 1.82, 0.06
    for i, (name, desc, accent) in enumerate(neta):
        col = i%4; row = i//4
        x = 0.35 + col*(cw+gap)
        y = 1.45 + row*(ch+gap)
        add_rect(slide, x, y, cw, ch, PANEL_BG)
        add_rect(slide, x, y, cw, 0.09, accent)
        add_rect(slide, x, y+ch-0.04, cw, 0.04, LIGHT_GRAY)
        add_text(slide, name, x+0.15, y+0.2,  cw-0.22, 0.6,  size=18, bold=True, color=DARK)
        add_text(slide, desc, x+0.15, y+0.92, cw-0.22, 0.82, size=12, color=CHARCOAL, wrap=True)

# ============================================================
# 8. 大将について
# ============================================================
def slide_chef(prs):
    slide = blank_slide(prs)
    header(slide, "大将について", "THE CHEF")
    items = [
        ("職人歴",   "25年以上"),
        ("修業先",   "銀座久兵衛系　老舗寿司店"),
        ("スタイル", "正統派江戸前鮨"),
        ("こだわり", "毎日自ら市場へ出向き、最高の素材を選び抜く"),
    ]
    y = 1.5
    for idx, (label, val) in enumerate(items):
        add_rect(slide, 0.4, y, 12.53, 1.1, PANEL_BG if idx%2==0 else WHITE)
        add_rect(slide, 0.4, y, 0.09, 1.1, GOLD)
        add_rect(slide, 0.4, y+1.07, 12.53, 0.03, LIGHT_GRAY)
        add_text(slide, label, 0.67, y+0.14, 2.2, 0.58, size=15, bold=True, color=GOLD)
        add_text(slide, val,   3.0,  y+0.14, 9.8, 0.75, size=20, bold=True, color=DARK, wrap=True)
        y += 1.13
    add_rect(slide, 0.4, y+0.18, 12.53, 1.28, GOLD_BG)
    add_rect(slide, 0.4, y+0.18, 12.53, 0.06, GOLD)
    add_rect(slide, 0.4, y+1.4,  12.53, 0.06, GOLD)
    add_text(slide,
             "「 握りにもつまみにも手間ひまかけ、緩急自在の味わいで\nお客様を江戸前の世界へご案内する 」",
             0.6, y+0.32, 12.1, 1.0,
             size=17, color=GOLD_DARK, align=PP_ALIGN.CENTER, italic=True)

# ============================================================
# 9. 口コミ・評判
# ============================================================
def slide_reviews(prs):
    slide = blank_slide(prs)
    header(slide, "口コミ・評判", "REVIEWS & REPUTATION")
    add_rect(slide, 0.4, 1.5, 3.5, 3.0, GOLD_BG)
    add_rect(slide, 0.4, 1.5, 3.5, 0.06, GOLD)
    add_text(slide, "食べログ", 0.5, 1.6, 3.3, 0.55, size=16, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "3.56",    0.5, 2.05, 3.3, 1.15, size=60, bold=True, color=GOLD_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, "534件のレビュー", 0.5, 3.3, 3.3, 0.45, size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
    points = [
        ("コストパフォーマンス", "1万円以下でも本格江戸前鮨が楽しめると高評価"),
        ("職人技",              "つまみと握り、両方に手間をかける丁寧な仕事"),
        ("雰囲気",              "個室チャージなしで高級感のある落ち着いた空間"),
        ("アクセス",            "上野御徒町駅から徒歩0分の抜群の立地"),
    ]
    y = 1.5
    for title, body in points:
        add_rect(slide, 4.1, y, 8.83, 1.0, PANEL_BG)
        add_rect(slide, 4.1, y, 0.09, 1.0, GOLD)
        add_rect(slide, 4.1, y+0.97, 8.83, 0.03, LIGHT_GRAY)
        add_text(slide, title, 4.38, y+0.1,  3.5, 0.48, size=15, bold=True, color=GOLD_DARK)
        add_text(slide, body,  4.38, y+0.55, 8.3, 0.42, size=13, color=DARK)
        y += 1.03
    add_rect(slide, 0.4, 5.6, 12.53, 0.06, GOLD)
    add_text(slide, "お客様の声", 0.4, 5.72, 5, 0.5, size=16, bold=True, color=GOLD_DARK)
    reviews = [
        "「1万円以下でこれだけの鮨が食べられるのは奇跡。またリピートします！」",
        "「のどぐろと炙り大トロが絶品。職人さんとの会話も楽しかった。」",
        "「上野でこんな本格派があるとは。接待にも使えるクオリティ。」",
    ]
    y = 6.28
    for rv in reviews:
        add_rect(slide, 0.4, y, 12.53, 0.38, PANEL_BG)
        add_text(slide, rv, 0.6, y+0.04, 12.1, 0.32, size=12, color=CHARCOAL, italic=True)
        y += 0.42

# ============================================================
# 10. まとめ・予約方法
# ============================================================
def slide_summary(prs):
    slide = blank_slide(prs)
    header(slide, "まとめ・予約方法", "SUMMARY & RESERVATION")
    add_rect(slide, 0.4, 1.5, 12.53, 1.25, GOLD_BG)
    add_rect(slide, 0.4, 1.5, 12.53, 0.06, GOLD)
    add_rect(slide, 0.4, 2.69, 12.53, 0.06, GOLD)
    add_text(slide,
             "銀座久兵衛の系譜を継ぐ正統派江戸前鮨を、上野で・リーズナブルに・個室チャージなしで楽しめる隠れた名店。\n"
             "職人歴25年の大将が自ら選び抜いた旬の素材を最高の技術で提供。コスパ・品質・雰囲気の三拍子が揃う。",
             0.6, 1.6, 12.1, 1.05, size=14, color=DARK, wrap=True)
    # 左：おすすめポイント
    add_rect(slide, 0.4, 2.85, 5.9, 4.3, PANEL_BG)
    add_rect(slide, 0.4, 2.85, 0.09, 4.3, GOLD)
    add_text(slide, "おすすめポイント", 0.67, 2.95, 5.5, 0.62, size=18, bold=True, color=GOLD_DARK)
    add_rect(slide, 0.5, 3.57, 5.7, 0.04, GOLD_LIGHT)
    recs = [
        "✦  6,000円〜本格江戸前鮨を体験できる",
        "✦  個室チャージ・サービス料なし",
        "✦  上野御徒町駅徒歩0分のアクセス",
        "✦  接待・記念日・デートに最適",
        "✦  四季折々の旬ネタが楽しめる",
        "✦  食べログ評価3.56（534件）",
    ]
    y = 3.68
    for r in recs:
        add_text(slide, r, 0.67, y, 5.5, 0.52, size=14, color=DARK)
        y += 0.52
    # 右：予約・連絡先
    add_rect(slide, 6.8, 2.85, 6.13, 4.3, PANEL_BG)
    add_rect(slide, 6.8, 2.85, 0.09, 4.3, GOLD)
    add_text(slide, "予約・連絡先", 7.07, 2.95, 5.8, 0.62, size=18, bold=True, color=GOLD_DARK)
    add_rect(slide, 6.9, 3.57, 5.9, 0.04, GOLD_LIGHT)
    res_info = [
        ("電話",      "03-6284-4731"),
        ("住所",      "東京都台東区上野4-4-5\nc-roadビル12F"),
        ("アクセス",  "上野御徒町駅A5出口 徒歩0分"),
        ("Instagram", "@ueno_sakae"),
        ("予約サイト","食べログ・一休.comレストラン"),
    ]
    y = 3.68
    for label, val in res_info:
        add_text(slide, label, 7.07, y, 2.0, 0.62, size=13, bold=True, color=GOLD)
        add_text(slide, val,   9.1,  y, 3.7, 0.72, size=14, color=DARK, wrap=True)
        add_rect(slide, 6.9, y+0.76, 5.9, 0.02, LIGHT_GRAY)
        y += 0.78
    add_rect(slide, 0, 7.2, 13.33, 0.3, GOLD_BG)
    add_rect(slide, 0, 7.2, 13.33, 0.05, GOLD)
    add_text(slide,
             "上野 榮（さかえ）　─　銀座久兵衛の系譜　正統派江戸前鮨　─　東京都台東区上野4-4-5 c-roadビル12F",
             0.3, 7.23, 12.7, 0.28, size=10, color=GOLD, align=PP_ALIGN.CENTER)

# ============================================================
# メイン
# ============================================================
prs = new_prs()
slide_cover(prs)
slide_overview(prs)
slide_concept(prs)
slide_atmosphere(prs)
slide_courses(prs)
slide_flow(prs)
slide_neta(prs)
slide_chef(prs)
slide_reviews(prs)
slide_summary(prs)

out = r"C:\Users\topge\OneDrive\ドキュメント\GitHub\claude_playground\docs\上野_榮_調査.pptx"
prs.save(out)
print(f"完了: {out}")
