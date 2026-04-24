from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── カラーパレット（白ベース・高級感） ───
C_BG    = RGBColor(0xFA, 0xF8, 0xF3)   # ウォームホワイト（背景）
C_PANEL = RGBColor(0xF0, 0xED, 0xE5)   # ライトクリーム（パネル）
C_CARD  = RGBColor(0xFF, 0xFF, 0xFF)   # 純白（カード）
C_NAVY  = RGBColor(0x14, 0x1C, 0x38)   # 深紺（見出しパネル）
C_NAVY2 = RGBColor(0x1E, 0x28, 0x50)   # やや明るい紺
C_GOLD  = RGBColor(0xB8, 0x90, 0x18)   # リッチゴールド（白背景向け）
C_GOLDB = RGBColor(0xD4, 0xAA, 0x2A)   # ゴールド明るめ
C_INK   = RGBColor(0x1A, 0x1A, 0x2E)   # 本文テキスト（ほぼ黒）
C_DGRAY = RGBColor(0x44, 0x44, 0x58)   # サブテキスト
C_MGRAY = RGBColor(0x88, 0x88, 0x9E)   # キャプション
C_LGRAY = RGBColor(0xD8, 0xD5, 0xCE)   # 区切り線・薄ボーダー
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN = RGBColor(0x14, 0x88, 0x48)
C_RED   = RGBColor(0xBE, 0x28, 0x28)
C_WARN  = RGBColor(0xBB, 0x70, 0x00)
C_LWARN = RGBColor(0xFF, 0xF3, 0xE0)   # 注意背景（薄オレンジ）
C_LRED  = RGBColor(0xFF, 0xED, 0xED)   # 注意背景（薄赤）

# 店舗ごとのアクセントカラー（白背景で映える深み色）
SHOP_COLORS = [
    RGBColor(0xC8, 0x40, 0x40),  # 01: バーミリオン
    RGBColor(0x0C, 0x88, 0x6A),  # 02: エメラルド
    RGBColor(0x7C, 0x3A, 0xCC),  # 03: パープル
    RGBColor(0xC8, 0x6E, 0x00),  # 04: アンバー
    RGBColor(0x18, 0x6A, 0xB4),  # 05: ロイヤルブルー
]

# 店舗ごとのアクセント薄色（カード背景用）
SHOP_LIGHT = [
    RGBColor(0xFF, 0xF0, 0xF0),
    RGBColor(0xE8, 0xF8, 0xF4),
    RGBColor(0xF3, 0xEC, 0xFF),
    RGBColor(0xFF, 0xF4, 0xE0),
    RGBColor(0xE8, 0xF3, 0xFF),
]

FONT_JP = "游ゴシック"

# ─── ヘルパー ───
def bg(slide):
    s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = C_BG; s.line.fill.background()

def rect(slide, x, y, w, h, fill=C_CARD, lc=None, lw=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc:
        s.line.color.rgb = lc
        s.line.width = Pt(lw or 1)
    else:
        s.line.fill.background()
    return s

def oval(slide, x, y, w, h, fill=C_GOLD):
    s = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, sz=18, bold=False, italic=False,
        color=C_INK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.italic = italic
    p.font.color.rgb = color; p.font.name = FONT_JP; p.alignment = align
    return tb

def stars(s):
    try:
        n = int(round(float(s)))
        return "★" * n + "☆" * (5 - n)
    except:
        return s

# ════════════════════════════════════════════════
# スライド 1: タイトル（白ベース）
# ════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)

# 左: 深紺ブロック（タイトルエリア）
rect(sl, 0, 0, 7.6, 7.5, C_NAVY)

# 装飾円（紺エリア内、右寄り）
oval(sl, 4.2, -1.2, 6.5, 6.5, C_NAVY2)
oval(sl, 5.5, -0.5, 4.5, 4.5, RGBColor(0x26, 0x32, 0x5E))

# ゴールドドット（紺エリア下部）
for i in range(6):
    oval(sl, 0.55 + i * 0.52, 6.55, 0.22, 0.22, C_GOLDB)

# 上部ゴールドバー
rect(sl, 0, 0, 7.6, 0.16, C_GOLDB)

# タグ（銀座線×浅草線）
rect(sl, 0.55, 0.52, 4.8, 0.52, RGBColor(0x26, 0x30, 0x58), C_GOLDB, 0.6)
txt(sl, "銀座線  ×  浅草線  沿線", 0.55, 0.52, 4.8, 0.52,
    sz=16, italic=True, color=C_GOLDB, align=PP_ALIGN.CENTER)

# メインタイトル
txt(sl, "厳選", 0.55, 1.22, 6.8, 1.10,
    sz=72, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
txt(sl, "寿司  5  選", 0.55, 2.22, 6.8, 1.20,
    sz=64, bold=True, color=C_GOLDB, align=PP_ALIGN.LEFT)

# ゴールドライン
rect(sl, 0.55, 3.55, 6.6, 0.07, C_GOLDB)

# サブコピー（白文字）
txt(sl, "予算 1 万強 / 人  ｜  男友達 3 人", 0.55, 3.72, 6.8, 0.55,
    sz=18, color=C_WHITE, align=PP_ALIGN.LEFT)
txt(sl, "肩ひじ張らない江戸前鮨", 0.55, 4.22, 6.8, 0.50,
    sz=18, color=RGBColor(0xB8, 0xC0, 0xD8), align=PP_ALIGN.LEFT)

# フィルター説明
rect(sl, 0.55, 5.0, 6.6, 0.72, RGBColor(0x22, 0x2C, 0x50), C_GOLDB, 0.6)
txt(sl, "食べログ × 口コミ件数 × メディア掲載\nでサクラを徹底排除",
    0.65, 5.05, 6.4, 0.62, sz=13, color=RGBColor(0xB8, 0xC4, 0xD8),
    align=PP_ALIGN.CENTER)

# 日付
txt(sl, "2026年 4月 調査", 5.5, 7.1, 1.9, 0.35,
    sz=11, italic=True, color=C_MGRAY, align=PP_ALIGN.RIGHT)

# 右: 白エリア装飾
rect(sl, 7.6, 0, 5.73, 7.5, C_BG)
rect(sl, 7.6, 0, 0.08, 7.5, C_GOLDB)  # 境界ゴールドライン

# 右エリアの装飾（大きな円）
oval(sl, 8.5, 0.5, 4.5, 4.5, C_PANEL)
oval(sl, 9.2, 1.2, 3.1, 3.1, C_CARD)

# 右エリア：5店の情報プレビュー
shops_preview = [
    (SHOP_COLORS[0], "浅草橋 鮨 うらおにかい", "浅草線"),
    (SHOP_COLORS[1], "銀座のみこ寿司",          "浅草線"),
    (SHOP_COLORS[2], "鮨結う遥",                "銀座線"),
    (SHOP_COLORS[3], "上野 榮",                 "銀座線"),
    (SHOP_COLORS[4], "すしのすけ",              "銀座線・浅草線"),
]
for i, (c, name, line) in enumerate(shops_preview):
    yy = 4.85 + i * 0.50
    rect(sl, 7.85, yy, 0.25, 0.32, c)
    txt(sl, name, 8.20, yy + 0.02, 3.5, 0.30, sz=13, bold=True, color=C_INK)
    txt(sl, line, 11.62, yy + 0.04, 1.45, 0.26, sz=11, italic=True,
        color=C_MGRAY, align=PP_ALIGN.RIGHT)

txt(sl, "PICK UP  5  RESTAURANTS", 7.85, 4.40, 5.2, 0.40,
    sz=13, bold=True, italic=True, color=C_MGRAY)


# ════════════════════════════════════════════════
# スライド 2: 選定基準（白ベース）
# ════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)

# ヘッダーバー（深紺）
rect(sl, 0, 0, 13.33, 1.05, C_NAVY)
rect(sl, 0, 0, 13.33, 0.10, C_GOLDB)  # 上部ゴールドライン
txt(sl, "選定基準  ─  サクラを見抜く 4 つのフィルター",
    0.45, 0.22, 12.0, 0.72, sz=26, bold=True, color=C_WHITE)

criteria = [
    (SHOP_COLORS[0], "①", "食べログ評価 3.5 以上",
     "AIによる不正検知システムあり。業者による操作レビューは定期的に削除される。3.5超えは本物の実力の証。"),
    (SHOP_COLORS[1], "②", "口コミ件数 100 件以上",
     "件数が多いほど1件のサクラが全体に与える影響が小さくなる。今回の採用店は最低でも255件。"),
    (SHOP_COLORS[2], "③", "複数の独立メディアに掲載",
     "東京カレンダー・ヒトサラ・タイムアウト東京などの編集部が取材・選定した店を優先。"),
    (SHOP_COLORS[3], "④", "実績・ブランド背景の確認",
     "老舗の暖簾分け、超人気店の姉妹店など、バックグラウンドが独立して検証できる店のみ採用。"),
]

for i, (color, num, title, desc) in enumerate(criteria):
    yy = 1.18 + i * 1.55
    # カード（白）
    rect(sl, 0.4, yy, 12.53, 1.32, C_CARD, C_LGRAY, 0.5)
    # 左アクセントバー
    rect(sl, 0.4, yy, 0.13, 1.32, color)
    # 番号サークル
    oval(sl, 0.68, yy + 0.30, 0.72, 0.72, color)
    txt(sl, num, 0.68, yy + 0.30, 0.72, 0.72,
        sz=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # タイトル
    txt(sl, title, 1.60, yy + 0.08, 10.9, 0.55,
        sz=20, bold=True, color=C_INK)
    # 説明
    txt(sl, desc, 1.60, yy + 0.65, 10.9, 0.58,
        sz=14, color=C_DGRAY)

# ════════════════════════════════════════════════
# 店舗スライド生成
# ════════════════════════════════════════════════
def shop_slide(idx, name, station, line, score, reviews, budget,
               trust, trust_note, points, caution, media):
    ac  = SHOP_COLORS[idx]
    alc = SHOP_LIGHT[idx]

    hex_ac = str(ac)
    r_ac = int(hex_ac[0:2], 16)
    g_ac = int(hex_ac[2:4], 16)
    b_ac = int(hex_ac[4:6], 16)
    ac_wm = RGBColor(
        min(255, int(r_ac * 0.12 + 0xFA * 0.88)),
        min(255, int(g_ac * 0.12 + 0xF8 * 0.88)),
        min(255, int(b_ac * 0.12 + 0xF3 * 0.88)),
    )

    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)

    # 大型ウォーターマーク番号（白背景上で極薄）
    txt(sl, f"0{idx+1}", -0.2, -0.3, 5.0, 4.0,
        sz=200, bold=True, color=ac_wm, align=PP_ALIGN.LEFT)

    # 上部アクセントバー（太め）
    rect(sl, 0, 0, 13.33, 0.22, ac)

    # ヘッダーパネル（深紺）
    rect(sl, 0, 0.22, 13.33, 1.18, C_NAVY)

    # 店名（白文字）
    txt(sl, name, 0.5, 0.28, 10.3, 0.88, sz=30, bold=True, color=C_WHITE)

    # 駅名 ＋ 路線
    oval(sl, 0.50, 1.09, 0.22, 0.22, ac)
    txt(sl, f"{station}  ／  {line}", 0.82, 1.06, 9.5, 0.35,
        sz=15, color=RGBColor(0xB8, 0xC8, 0xE0))

    # Noバッジ（右上）
    rect(sl, 11.50, 0.30, 1.60, 0.76, ac)
    txt(sl, f"No.{idx+1:02d}", 11.50, 0.30, 1.60, 0.76,
        sz=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ════ 左カード（白、カラーボーダー）════
    LX, LW = 0.28, 3.85
    rect(sl, LX, 1.58, LW, 5.72, C_CARD, ac, 1.2)
    rect(sl, LX, 1.58, 0.10, 5.72, ac)

    # スコアエリア薄色背景
    rect(sl, LX + 0.10, 1.58, LW - 0.10, 1.85, alc)

    txt(sl, "食べログ", LX, 1.68, LW, 0.36,
        sz=11, color=C_MGRAY, align=PP_ALIGN.CENTER)
    txt(sl, score, LX, 1.98, LW, 0.85,
        sz=46, bold=True, color=ac, align=PP_ALIGN.CENTER)
    txt(sl, stars(score), LX, 2.77, LW, 0.40,
        sz=19, color=ac, align=PP_ALIGN.CENTER)
    txt(sl, f"口コミ  {reviews} 件", LX, 3.11, LW, 0.34,
        sz=12, color=C_DGRAY, align=PP_ALIGN.CENTER)

    rect(sl, LX + 0.22, 3.50, LW - 0.44, 0.04, C_LGRAY)

    txt(sl, "予算（夜）", LX, 3.60, LW, 0.34,
        sz=11, color=C_MGRAY, align=PP_ALIGN.CENTER)
    txt(sl, budget, LX, 3.90, LW, 0.58,
        sz=16, bold=True, color=C_INK, align=PP_ALIGN.CENTER)

    rect(sl, LX + 0.22, 4.52, LW - 0.44, 0.04, C_LGRAY)

    txt(sl, "サクラ安心度", LX, 4.60, LW, 0.34,
        sz=11, color=C_MGRAY, align=PP_ALIGN.CENTER)
    tc = C_GREEN if trust >= 4 else (C_WARN if trust == 3 else C_RED)
    txt(sl, "●" * trust + "○" * (5 - trust), LX, 4.90, LW, 0.46,
        sz=22, bold=True, color=tc, align=PP_ALIGN.CENTER)
    txt(sl, trust_note, LX, 5.30, LW, 0.35,
        sz=10, color=C_DGRAY, align=PP_ALIGN.CENTER)

    rect(sl, LX + 0.22, 5.70, LW - 0.44, 0.04, C_LGRAY)

    txt(sl, "掲載メディア", LX, 5.78, LW, 0.32,
        sz=11, color=C_MGRAY, align=PP_ALIGN.CENTER)
    txt(sl, media, LX, 6.08, LW, 0.95,
        sz=11, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # ════ 右セクション ════
    RX, RW = 4.38, 8.70

    # ポイントヘッダー（アクセント帯）
    rect(sl, RX, 1.58, RW, 0.46, ac)
    txt(sl, "  おすすめポイント", RX, 1.58, RW, 0.46,
        sz=16, bold=True, color=C_WHITE)

    # ポイント5項目（白カード）
    y_pt = 2.10
    for pt in points:
        rect(sl, RX, y_pt, RW, 0.60, C_CARD, C_LGRAY, 0.5)
        rect(sl, RX, y_pt, 0.10, 0.60, ac)
        txt(sl, pt, RX + 0.22, y_pt + 0.07, RW - 0.32, 0.48,
            sz=13, color=C_INK)
        y_pt += 0.65

    # 注意ボックス（薄オレンジ）
    if caution:
        y_c = y_pt + 0.06
        rect(sl, RX, y_c, RW, 1.16, C_LWARN, C_WARN, 0.6)
        rect(sl, RX, y_c, 0.10, 1.16, C_WARN)
        txt(sl, "⚠  注意点", RX + 0.22, y_c + 0.06, RW - 0.32, 0.36,
            sz=13, bold=True, color=C_WARN)
        txt(sl, caution, RX + 0.22, y_c + 0.44, RW - 0.32, 0.66,
            sz=12, color=C_DGRAY)


# ════════════════════════════════════════════════
# スライド 3〜7: 各店舗
# ════════════════════════════════════════════════
shop_slide(
    idx=0,
    name="浅草橋 鮨 うらおにかい",
    station="浅草橋駅 徒歩2分",
    line="都営浅草線",
    score="3.53", reviews="255",
    budget="10,000〜14,999円",
    trust=4, trust_note="255件・複数メディア掲載",
    points=[
        "裏路地の「秘密基地」風外観 → 入った瞬間テンション爆上がり",
        "くずし鮨コース：15貫＋小皿3品＋1ドリンク付き",
        "若手職人に活躍の場を与えるため2019年に誕生した革新系",
        "海老天ノリ巻き・季節おまかせなど遊び心が満載",
        "タイムアウト東京・ヒトサラ 両メディアに掲載",
    ],
    caution="飲み物を追加すると1.5万を超える場合あり。事前にコース内容と込み料金を確認推奨。",
    media="タイムアウト東京 / ヒトサラ / Retty",
)

shop_slide(
    idx=1,
    name="銀座のみこ寿司",
    station="東銀座駅 徒歩4分",
    line="都営浅草線",
    score="3.57", reviews="433",
    budget="8,000〜9,999円",
    trust=4, trust_note="433件（5店中最多）",
    points=[
        "口コミ433件は本日紹介する5店の中でダントツ1位",
        "毎朝市場から直仕入れの旬ネタをリーズナブルに提供",
        "「ウニパラダイス」— 複数産地のウニを食べ比べられる名物",
        "カウンター10席のみ。職人との距離が近く会話が弾む",
        "予算1万円以内に収まる可能性が高い（5店で最も安い帯）",
    ],
    caution="ウニ専門色が強め。ウニが苦手なメンバーがいる場合は要確認。ランチの食べ放題と混同注意。",
    media="ヒトサラマガジン / 食べログ",
)

shop_slide(
    idx=2,
    name="鮨結う遥  （すしゆうはるか）",
    station="末広町駅 徒歩3分",
    line="東京メトロ銀座線",
    score="3.65", reviews="新店につき蓄積中",
    budget="13,200円（飲み放題込み）",
    trust=3, trust_note="新店・親ブランドで信頼補完",
    points=[
        "食べログ 3.65 ── 今回紹介する5店の中で最高スコア",
        "超人気「鮨結う翼」（恵比寿）の姉妹店。職人の質は折り紙付き",
        "飲み放題込み13,200円 → 酒好き3人なら実質的に最安値級",
        "若大将の軽快なトークでエンタメ感があり、男友達3人向け",
        "東京カレンダー系メディア「グルカレ」で特集記事掲載",
    ],
    caution="2024年12月開店の新店。口コミ件数が少なく評価の母数が小さい。予約は数ヶ月先まで埋まりやすい。",
    media="グルカレ（東京カレンダー系）/ 一休.com",
)

shop_slide(
    idx=3,
    name="上野 榮",
    station="上野広小路駅 すぐ",
    line="東京メトロ銀座線",
    score="3.56", reviews="319",
    budget="9,000〜13,200円",
    trust=4, trust_note="319件・一休掲載・評価安定",
    points=[
        "雑居ビル12階の隠れ家。上野の下町感と好対照の洗練された空間",
        "319件の口コミで評価が安定 ── 長期的な実力店の証",
        "個室チャージなし・サービス料なしという良心的な会計",
        "おまかせ9,000円〜 と今回の中でも手頃なスタートライン",
        "赤酢シャリ・岩塩で食べるスタイルで鮨好きも唸る技",
    ],
    caution="完全予約制のおまかせコースのみ。飛び込みは不可。コースによっては13,200円になる場合あり。",
    media="一休.comレストラン / Retty",
)

shop_slide(
    idx=4,
    name="おまかせ寿司 すしのすけ",
    station="新橋駅 徒歩1分",
    line="銀座線・浅草線 両線利用可",
    score="3.52", reviews="21",
    budget="8,800円（おまかせ20品）",
    trust=2, trust_note="口コミ21件 → 要注意",
    points=[
        "新橋駅1分・ニュー新橋ビルB1F。5店の中でアクセス最強",
        "おまかせ20品 8,800円 ── 今回の5店でダントツ最安値",
        "BAR風のムーディな内装で「カジュアル感」が最も高い",
        "ウニ・イクラ・カニ乗せ小丼が名物料理として評判",
        "数ヶ月先まで予約困難 → 実際の来客に支えられた人気",
    ],
    caution="⚠ 口コミ数21件は5店中最少で評価の信頼性が最も低い。複数の個人ブログで裏取りは取れているが、来店前に直近の口コミを必ず確認すること。",
    media="note（ウニ王子）/ 私的標本ブログ / ヒトサラ",
)


# ════════════════════════════════════════════════
# スライド 8: 比較まとめ（白ベース・最寄り駅列）
# ════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)

# ヘッダーバー（深紺）
rect(sl, 0, 0, 13.33, 1.05, C_NAVY)
rect(sl, 0, 0, 13.33, 0.10, C_GOLDB)
txt(sl, "5 店  比較まとめ", 0.35, 0.18, 9.0, 0.76,
    sz=28, bold=True, color=C_WHITE)

# テーブル定義
headers    = ["店名", "最寄り駅", "食べログ", "予算/人", "安心度", "一言"]
col_widths = [2.72,   2.55,       1.50,        2.00,      1.20,     2.76]
col_x = [0.20]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

# ヘッダー行（ゴールド）
y_hdr = 1.12
rect(sl, 0.20, y_hdr, sum(col_widths), 0.50, C_GOLDB)
for h, cx, cw in zip(headers, col_x, col_widths):
    txt(sl, h, cx + 0.04, y_hdr + 0.05, cw - 0.08, 0.40,
        sz=13, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

rows_data = [
    ["浅草橋 鮨 うらおにかい", "浅草橋駅 徒歩2分",  "3.53 / 255件", "〜15,000円",     ("●●●●○", 4), "くずし鮨×秘密基地"],
    ["銀座のみこ寿司",          "東銀座駅 徒歩4分",  "3.57 / 433件", "〜9,999円",      ("●●●●○", 4), "ウニ特化・最多口コミ"],
    ["鮨結う遥",                "末広町駅 徒歩3分",  "3.65 / 新店",  "13,200円(飲放)", ("●●●○○", 3), "飲み放題込・最高点"],
    ["上野 榮",                 "上野広小路駅すぐ",  "3.56 / 319件", "9,000円〜",      ("●●●●○", 4), "安定実力・個室あり"],
    ["おまかせ寿司すしのすけ",  "新橋駅 徒歩1分",    "3.52 / 21件",  "8,800円",        ("●●○○○", 2), "最安値・口コミ要確認"],
]

row_h = 1.04
for ri, row in enumerate(rows_data):
    yy = y_hdr + 0.50 + ri * row_h
    fill_c = C_CARD if ri % 2 == 0 else C_PANEL
    rect(sl, 0.20, yy, sum(col_widths), row_h, fill_c, C_LGRAY, 0.3)
    # 店舗カラー左帯
    rect(sl, 0.20, yy, 0.12, row_h, SHOP_COLORS[ri])

    cells = row[:4] + [row[4][0]] + [row[5]]
    for ci, (cell, cx, cw) in enumerate(zip(cells, col_x, col_widths)):
        c = C_INK
        if ci == 2:
            c = C_GOLD
        elif ci == 4:
            n = row[4][1]
            c = C_GREEN if n >= 4 else (C_WARN if n == 3 else C_RED)
        txt(sl, cell, cx + 0.14, yy + 0.10, cw - 0.20, row_h - 0.20,
            sz=12, color=c, align=PP_ALIGN.CENTER)

txt(sl, "※ 情報は2026年4月時点。最新の予約状況・価格は各店公式サイト・食べログで必ずご確認ください。",
    0.20, 7.12, 12.9, 0.35, sz=10, color=C_MGRAY)


# ════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════
out = r"C:\Users\topge\OneDrive\ドキュメント\GitHub\claude_playground\games\sushi_ginza_asakusa.pptx"
prs.save(out)
print(f"保存完了: {out}")
